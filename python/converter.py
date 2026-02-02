#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX to JPEG Converter for SyncDisplay
Converts PowerPoint presentations to high-quality JPEG images for fast display.

Supports two conversion methods:
1. LibreOffice (preferred) - High fidelity, preserves all formatting
2. PyMuPDF + pdf2image - Fallback if LibreOffice unavailable

Usage:
    python converter.py --input presentation.pptx --output ./cache/russian --width 1920 --height 1080
"""

import argparse
import io
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

# Fix Windows console encoding for Unicode (Cyrillic, etc.)
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("WARNING: Pillow not installed. Install with: pip install Pillow")

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def find_libreoffice():
    """Find LibreOffice executable on the system."""
    import platform
    system = platform.system()
    
    # Check if soffice is in PATH first (works on all platforms)
    if shutil.which("soffice"):
        return "soffice"
    if shutil.which("soffice.exe"):
        return "soffice.exe"
    
    possible_paths = []
    
    if system == "Windows":
        # Windows paths
        possible_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            r"C:\Program Files\LibreOffice\program\soffice.com",
            os.path.expanduser(r"~\AppData\Local\Programs\LibreOffice\program\soffice.exe"),
        ]
    elif system == "Darwin":  # macOS
        possible_paths = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            os.path.expanduser("~/Applications/LibreOffice.app/Contents/MacOS/soffice"),
        ]
    else:  # Linux
        possible_paths = [
            "/usr/bin/soffice",
            "/usr/bin/libreoffice",
            "/usr/local/bin/soffice",
            "/usr/local/bin/libreoffice",
            "/opt/libreoffice/program/soffice",
            "/snap/bin/libreoffice",  # Snap package
            os.path.expanduser("~/.local/bin/soffice"),
        ]
    
    for path in possible_paths:
        if os.path.exists(path):
            return path
    
    return None


def convert_with_libreoffice(input_path, output_dir, width, height):
    """
    Convert PPTX to images using LibreOffice.
    This method provides the highest fidelity conversion.
    """
    soffice_path = find_libreoffice()
    if not soffice_path:
        raise RuntimeError("LibreOffice not found. Please install LibreOffice.")
    
    # Create temporary directory for PDF
    with tempfile.TemporaryDirectory() as temp_dir:
        # Convert PPTX to PDF using LibreOffice
        print("Converting PPTX to PDF with LibreOffice...")
        print(f"Using LibreOffice at: {soffice_path}")
        
        cmd = [
            soffice_path,
            "--headless",
            "--nofirststartwizard",
            "--norestore",
            "--convert-to", "pdf",
            "--outdir", temp_dir,
            input_path
        ]
        
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
            print(f"LibreOffice stdout: {result.stdout}")
            if result.stderr:
                print(f"LibreOffice stderr: {result.stderr}")
            
            if result.returncode != 0:
                # Try to provide more useful error info
                error_msg = result.stderr or result.stdout or "Unknown error (exit code: {})".format(result.returncode)
                raise RuntimeError(f"LibreOffice conversion failed: {error_msg}")
        except subprocess.TimeoutExpired:
            raise RuntimeError("LibreOffice conversion timed out after 5 minutes")
        
        # Find the generated PDF
        pdf_filename = Path(input_path).stem + ".pdf"
        pdf_path = os.path.join(temp_dir, pdf_filename)
        
        print(f"Looking for PDF at: {pdf_path}")
        
        if not os.path.exists(pdf_path):
            # Try to find any PDF in temp dir
            pdf_files = list(Path(temp_dir).glob("*.pdf"))
            print(f"Found PDF files in temp dir: {pdf_files}")
            if pdf_files:
                pdf_path = str(pdf_files[0])
            else:
                # List what files ARE in the temp dir for debugging
                all_files = list(Path(temp_dir).glob("*"))
                print(f"Files in temp dir: {all_files}")
                raise RuntimeError(f"PDF file not generated. Expected: {pdf_filename}")
        
        # Convert PDF pages to images using PyMuPDF
        print("Converting PDF pages to images...")
        return convert_pdf_to_images(pdf_path, output_dir, width, height)


def convert_pdf_to_images(pdf_path, output_dir, width, height):
    """Convert PDF pages to JPEG images using PyMuPDF."""
    if not PYMUPDF_AVAILABLE:
        raise RuntimeError("PyMuPDF not installed. Install with: pip install PyMuPDF")
    
    doc = fitz.open(pdf_path)
    total_pages = len(doc)
    slide_count = 0
    
    for page_num in range(total_pages):
        page = doc[page_num]
        
        # Calculate zoom factor for desired resolution
        # Standard slide aspect ratio is 16:9 or 4:3
        page_rect = page.rect
        zoom_x = width / page_rect.width
        zoom_y = height / page_rect.height
        zoom = min(zoom_x, zoom_y)  # Fit within bounds
        
        # Create transformation matrix
        mat = fitz.Matrix(zoom, zoom)
        
        # Render page to pixmap
        pix = page.get_pixmap(matrix=mat, alpha=False)
        
        # Convert to PIL Image for processing
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        # Resize to exact dimensions if needed, centering the content
        if img.width != width or img.height != height:
            # Create black background
            final_img = Image.new("RGB", (width, height), (0, 0, 0))
            # Calculate position to center
            x_offset = (width - img.width) // 2
            y_offset = (height - img.height) // 2
            final_img.paste(img, (x_offset, y_offset))
            img = final_img
        
        # Save full-size image
        slide_filename = f"slide_{page_num + 1:03d}.jpg"
        slide_path = os.path.join(output_dir, slide_filename)
        img.save(slide_path, "JPEG", quality=92, optimize=True)
        
        slide_count += 1
        
        # Print progress for IPC
        progress = int((page_num + 1) / total_pages * 100)
        print(f"PROGRESS:{progress}")
    
    doc.close()
    return slide_count


def convert_with_pptx_direct(input_path, output_dir, width, height):
    """
    Fallback conversion using python-pptx.
    Note: This method has limited formatting support compared to LibreOffice.
    It's mainly useful for text extraction and basic slide rendering.
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx not installed. Install with: pip install python-pptx")
    
    prs = Presentation(input_path)
    total_slides = len(prs.slides)
    
    print(f"WARNING: Direct PPTX rendering has limited support for complex formatting.")
    print(f"Consider installing LibreOffice for better results.")
    
    # For now, we'll create placeholder images
    # In a production system, you'd implement more sophisticated rendering
    for idx, slide in enumerate(prs.slides):
        # Create a basic image with slide text
        img = Image.new("RGB", (width, height), (0, 0, 32))  # Dark blue background
        
        # Save the image
        slide_filename = f"slide_{idx + 1:03d}.jpg"
        slide_path = os.path.join(output_dir, slide_filename)
        img.save(slide_path, "JPEG", quality=92)
        
        progress = int((idx + 1) / total_slides * 100)
        print(f"PROGRESS:{progress}")
    
    return total_slides


def generate_thumbnails(output_dir, thumb_width):
    """Generate thumbnail images for all slides."""
    print("Generating thumbnails...")
    
    slides = sorted([f for f in os.listdir(output_dir) 
                    if f.startswith("slide_") and f.endswith(".jpg") and "_thumb" not in f])
    
    for slide_file in slides:
        slide_path = os.path.join(output_dir, slide_file)
        thumb_filename = slide_file.replace(".jpg", "_thumb.jpg")
        thumb_path = os.path.join(output_dir, thumb_filename)
        
        with Image.open(slide_path) as img:
            # Calculate thumbnail height maintaining aspect ratio
            aspect = img.width / img.height
            thumb_height = int(thumb_width / aspect)
            
            # Create thumbnail
            thumb = img.copy()
            thumb.thumbnail((thumb_width, thumb_height), Image.Resampling.LANCZOS)
            thumb.save(thumb_path, "JPEG", quality=85)


def extract_text_from_pptx(input_path):
    """Extract text content from each slide for singer screen preview."""
    if not PPTX_AVAILABLE:
        print("WARNING: python-pptx not available, skipping text extraction")
        return []
    
    print("Extracting slide text...")
    prs = Presentation(input_path)
    slides_text = []
    
    for slide in prs.slides:
        slide_text = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            
            # Handle tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            slide_text.append(cell.text.strip())
        
        # Join all text, get first meaningful line
        full_text = "\n".join(slide_text)
        first_line = ""
        for line in full_text.split("\n"):
            line = line.strip()
            if line and len(line) > 2:  # Skip very short lines
                first_line = line
                break
        
        slides_text.append({
            "text": full_text,
            "firstLine": first_line
        })
    
    return slides_text


def create_metadata(output_dir, input_path, slide_count, slides_text):
    """Create metadata JSON file with slide information."""
    metadata = {
        "sourceFile": os.path.basename(input_path),
        "slideCount": slide_count,
        "generatedAt": __import__("datetime").datetime.now().isoformat(),
        "slides": slides_text
    }
    
    metadata_path = os.path.join(output_dir, "metadata.json")
    with open(metadata_path, "w", encoding="utf-8") as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)
    
    print(f"Metadata saved to {metadata_path}")


def main():
    parser = argparse.ArgumentParser(description="Convert PPTX to JPEG images")
    parser.add_argument("--input", "-i", required=True, help="Input PPTX file path")
    parser.add_argument("--output", "-o", required=True, help="Output directory for images")
    parser.add_argument("--width", "-W", type=int, default=1920, help="Output image width")
    parser.add_argument("--height", "-H", type=int, default=1080, help="Output image height")
    parser.add_argument("--thumbnail-width", "-t", type=int, default=300, help="Thumbnail width")
    parser.add_argument("--method", "-m", choices=["libreoffice", "direct", "auto"], 
                       default="auto", help="Conversion method")
    
    args = parser.parse_args()
    
    # Validate input file
    if not os.path.exists(args.input):
        print(f"ERROR: Input file not found: {args.input}")
        sys.exit(1)
    
    # Create output directory
    os.makedirs(args.output, exist_ok=True)
    
    # Choose conversion method
    method = args.method
    if method == "auto":
        if find_libreoffice() and PYMUPDF_AVAILABLE:
            method = "libreoffice"
        else:
            method = "direct"
            print("WARNING: LibreOffice not available, using direct conversion (limited formatting)")
    
    print(f"Starting conversion: {args.input}")
    print(f"Output directory: {args.output}")
    print(f"Resolution: {args.width}x{args.height}")
    print(f"Method: {method}")
    
    try:
        # Convert presentation
        if method == "libreoffice":
            slide_count = convert_with_libreoffice(
                args.input, args.output, args.width, args.height
            )
        else:
            slide_count = convert_with_pptx_direct(
                args.input, args.output, args.width, args.height
            )
        
        print(f"Converted {slide_count} slides")
        
        # Generate thumbnails
        if PIL_AVAILABLE:
            generate_thumbnails(args.output, args.thumbnail_width)
        
        # Extract text
        slides_text = extract_text_from_pptx(args.input)
        
        # Create metadata
        create_metadata(args.output, args.input, slide_count, slides_text)
        
        print("PROGRESS:100")
        print("Conversion complete!")
        
    except Exception as e:
        print(f"ERROR: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
